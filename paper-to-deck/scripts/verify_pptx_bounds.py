"""
verify_pptx_bounds.py — pre-declaration gate for .pptx builds.

Enforces three pptx-gotchas invariants *before* the user sees the deck:

  §7  No shape may extend outside the slide rectangle [0, slide_w] × [0, slide_h].
  §8  No textbox's first-run font size may exceed the available box height
      (size_pt / 72 > h - vertical margins).
  §9  Slides tagged as native-table (in outline.md) must NOT contain a Picture
      shape that looks like a table raster (tbl-*.png or similar).

Exit codes:
  0 = all gates passed
  1 = one or more violations found
  2 = usage / file error

Usage:
  python scripts/verify_pptx_bounds.py <deck.pptx>
    [--outline <outline.md>]   (optional, needed for §9 check)
    [--tolerance 0.02]         (inches; slack for floating-point rounding)

Notes:
- Does NOT render the deck (no LibreOffice / Chrome dependency).
- Does NOT catch every PPTX bug — it catches the class of bugs that produce
  off-slide or clipped content, which are the easiest to miss and the most
  embarrassing to ship.
- Windows cp950 safety: all stdout is ASCII.
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.stdout.reconfigure(encoding="utf-8", errors="replace")


EMU_PER_INCH = 914400


def emu_to_in(emu: int) -> float:
    return emu / EMU_PER_INCH


def walk_shapes(shape):
    """Yield shape + any descendants inside a group shape."""
    yield shape
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            yield from walk_shapes(sub)


def check_bounds(prs, tolerance: float) -> list[str]:
    """§7: shape bounding box must be inside slide rectangle."""
    sw = emu_to_in(prs.slide_width)
    sh = emu_to_in(prs.slide_height)
    violations = []
    for si, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            for sh_ in walk_shapes(shape):
                # some shapes (Picture placeholder on layout) may not have all
                # positions populated — skip those
                if sh_.left is None or sh_.top is None:
                    continue
                x = emu_to_in(sh_.left)
                y = emu_to_in(sh_.top)
                w = emu_to_in(sh_.width) if sh_.width else 0.0
                h = emu_to_in(sh_.height) if sh_.height else 0.0
                if x < -tolerance:
                    violations.append(
                        f"[slide {si:02d}] §7 shape {sh_.shape_id} starts at x={x:.3f} (< 0)"
                    )
                if y < -tolerance:
                    violations.append(
                        f"[slide {si:02d}] §7 shape {sh_.shape_id} starts at y={y:.3f} (< 0)"
                    )
                if x + w > sw + tolerance:
                    violations.append(
                        f"[slide {si:02d}] §7 shape {sh_.shape_id} extends to x+w={x+w:.3f} "
                        f"(slide width = {sw:.3f})"
                    )
                if y + h > sh + tolerance:
                    violations.append(
                        f"[slide {si:02d}] §7 shape {sh_.shape_id} extends to y+h={y+h:.3f} "
                        f"(slide height = {sh:.3f})"
                    )
    return violations


def check_hero_font(prs) -> list[str]:
    """§8: for any textbox whose first run has size_pt, make sure
    size_pt / 72 <= box_height - approx vertical margin."""
    violations = []
    # default text-frame vertical margin in python-pptx is 45 720 EMU top/bottom
    # (~0.05 in). If user sets 0, we still allow 0.05 slack.
    default_margin_in = 0.05
    for si, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.height is None:
                continue
            tf = shape.text_frame
            if not tf.paragraphs or not tf.paragraphs[0].runs:
                continue
            first_run = tf.paragraphs[0].runs[0]
            if first_run.font.size is None:
                continue
            size_pt = first_run.font.size.pt
            h_in = emu_to_in(shape.height)
            # account for top+bottom margin (use tf margins if set, else default)
            top_m_in = emu_to_in(tf.margin_top or 0) or default_margin_in
            bot_m_in = emu_to_in(tf.margin_bottom or 0) or default_margin_in
            usable_in = h_in - top_m_in - bot_m_in
            needed_in = size_pt / 72.0
            # only flag when the gap is substantial (>10 % short) — fonts have
            # some slack from line-height and renderer behaviour
            if needed_in > usable_in * 1.10:
                violations.append(
                    f"[slide {si:02d}] §8 textbox {shape.shape_id}: font {size_pt:.0f}pt "
                    f"(~{needed_in:.2f}in) exceeds usable height {usable_in:.2f}in "
                    f"(box height {h_in:.2f}in); text may clip"
                )
    return violations


def check_raster_tables(prs, outline_path: Path | None) -> list[str]:
    """§9: slides tagged V5 / native-table in outline.md MUST contain a native
    table shape (python-pptx GraphicFrame with .has_table). If not, the builder
    took the raster shortcut — flag.

    Rationale (design note): an earlier version of this check sniffed
    Picture.image.filename for "tbl-*.png" to decide if the picture was a table
    raster. This is unreliable — python-pptx renames embedded images to
    `imageN.png` inside the package zip, so the original filename is lost. The
    robust formulation is inverse: a V5 slide should HAVE a native table, and
    if it doesn't, something went wrong regardless of what the picture claims
    to be.
    """
    if outline_path is None or not outline_path.exists():
        return []  # can't evaluate without outline
    text = outline_path.read_text(encoding="utf-8", errors="replace")
    # find slides tagged V5 or explicitly "native table" in the outline.
    # Split on ANY `\n## ` header (not only `## Slide `) so trailing sections
    # like `## Caveats flagged for build step` or `## Requested edits` don't
    # get lumped into the last slide's block and trigger false positives.
    native_table_slides: set[int] = set()
    for block in re.split(r"\n## ", text):
        m = re.match(r"Slide\s+(\d+)", block)
        if not m:
            continue
        idx = int(m.group(1))
        # Match V5 as a word-boundary token (so "V50" etc. won't false-positive).
        if re.search(r"\bV5\b", block) or re.search(
                r"native\s+(?:editable\s+)?(?:PP?T\s+)?table", block, re.IGNORECASE):
            native_table_slides.add(idx)
    if not native_table_slides:
        return []
    violations = []
    for si, slide in enumerate(prs.slides, start=1):
        if si not in native_table_slides:
            continue
        has_native_table = any(
            getattr(s, "has_table", False) for s in slide.shapes
        )
        if has_native_table:
            continue  # satisfied — builder did the right thing
        pictures = [
            s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        if pictures:
            violations.append(
                f"[slide {si:02d}] §9 outline marks this slide V5 / native-table "
                f"but deck contains {len(pictures)} picture shape(s) and no "
                f"native table — raster shortcut taken"
            )
        else:
            violations.append(
                f"[slide {si:02d}] §9 outline marks this slide V5 / native-table "
                f"but deck contains no native table shape"
            )
    return violations


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx_path", type=Path)
    ap.add_argument("--outline", type=Path, default=None,
                    help="Path to outline.md (enables §9 check)")
    ap.add_argument("--tolerance", type=float, default=0.02,
                    help="Coordinate tolerance in inches (default 0.02)")
    args = ap.parse_args()

    if not args.pptx_path.exists():
        print(f"[ERR] {args.pptx_path} not found", file=sys.stderr)
        sys.exit(2)

    prs = Presentation(str(args.pptx_path))
    print(f"[INFO] verifying {args.pptx_path.name}  "
          f"slides={len(prs.slides)}  "
          f"canvas={emu_to_in(prs.slide_width):.2f}x{emu_to_in(prs.slide_height):.2f}in")

    all_violations: list[str] = []
    all_violations += check_bounds(prs, args.tolerance)
    all_violations += check_hero_font(prs)
    all_violations += check_raster_tables(prs, args.outline)

    if not all_violations:
        print("[OK] all gates passed (gotchas §7 / §8 / §9).")
        sys.exit(0)

    print()
    print(f"[FAIL] {len(all_violations)} violation(s):")
    for v in all_violations:
        print("  " + v)
    print()
    print("Fix the builder script (do not hand-patch in PowerPoint — changes are "
          "lost on regeneration). See paper-to-deck/references/pptx-gotchas.md.")
    sys.exit(1)


if __name__ == "__main__":
    main()
